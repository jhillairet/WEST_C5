# -*- coding: utf-8 -*-
"""
Generate a pptx file which concatenates all pulse snapshot figures

"""
#%%
from pptx import Presentation
from pptx.util import Inches
from glob import glob
from pathlib import Path
#%%
# define session name 
session_dir = './snapshots/2020-12-18_W-sources'
session_title = '18/12/2020 Impact of Boron-Nitride tiles on RF-induced W sources during plasma on WEST'

session_dir = './snapshots/2021-01-20_HC'
session_title = '20/01/2021 HC'

session_dir = './snapshots/2021-01-26_HC/'
session_title = '26/01/2021 High Confinement'

session_dir = './snapshots/2021-01-27_HC/'
session_title = '27/01/2021 High Confinement'

session_dir = './snapshots/2021-01-21_W-source Part2'
session_title = '21/01/2021 - Impact of Boron-Nitride tiles on RF-induced W sources during plasma on WEST - Part 2'

#%%
session_dir = Path(session_dir)
# list all .png files in a session directory
snapshot_filenames = glob(session_dir.as_posix()+'/*.png')

# create a presentation with one slide per snapshot figure 
prs = Presentation()

# first slide: title
slide = prs.slides.add_slide(prs.slide_layouts[0])
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = session_title

bullet_slide_layout = prs.slide_layouts[1]
for snapshot_filename in snapshot_filenames:
    pulse = snapshot_filename.split('WEST_')[-1].split('_')[0]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = f'#{pulse}'
    
    tf = body_shape.text_frame
    tf.text = 'Comments'

    # p = tf.add_paragraph()
    # p.text = 'Use _TextFrame.text for first bullet'
    # p.level = 1

    # p = tf.add_paragraph()
    # p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    # p.level = 2

    left = Inches(2)
    top = Inches(3)
    
    height = Inches(4.5)
    pic = slide.shapes.add_picture(snapshot_filename, left, top, height=height)

#%%
filename = str((session_dir/'session_summary.pptx').as_posix())
prs.save(filename)
